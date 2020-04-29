import os
from pptx import Presentation 
from pptx.util import Inches

img_path = os.getcwd() + '\\6.RdeT-6-3-Annealing.svg'

prs = Presentation() 
blank_slide_layout = prs.slide_layouts[6] 
slide = prs.slides.add_slide(blank_slide_layout)
print(os.getcwd())
print(img_path)
left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top) 
prs.save('test.pptx')