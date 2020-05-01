import glob
import os
import pptx
from pptx.util import Inches
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM

f = glob.glob(r'*.svg') # Grab the file name of a file type in the local folder

for img_path in f:
    print(img_path)
    drawing = svg2rlg(img_path)
    print(drawing)
    renderPM.drawToFile(drawing, img_path + '.png', fmt='PNG')
print(img_path)

pptFile = pptx.Presentation()

picFiles = [fn for fn in os.listdir() if fn.endswith('.png')]

# 按图片编号顺序导入
for fn in picFiles:
    slide = pptFile.slides.add_slide(pptFile.slide_layouts[1])

    # 为PPTX文件当前幻灯片中第一个文本框设置文字，本文代码中可忽略
    slide.shapes.placeholders[0].text = fn[:fn.rindex('.')]

    # 导入并为当前幻灯片添加图片，起始位置和尺寸可修改
    slide.shapes.add_picture(fn, Inches(0), Inches(0), Inches(10), Inches(7.5))
   
pptFile.save('<filename>.pptx')