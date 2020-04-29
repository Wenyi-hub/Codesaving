import os
os.chdir("<filename.filetype>")
from pptx import Presentation
from pptx.util import Cm

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(bullet_slide_layout)

left = top = Cm(3)
height = Cm(5)
pic = slide.shapes.add_picture("demo.png",left,top,height = height)

pic = slide.shapes.add_picture("demo.png",left,top,height)

prs.save('<filename.filetype>')